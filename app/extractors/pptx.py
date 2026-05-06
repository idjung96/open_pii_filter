"""PPTX text extraction (Phase 4b).

Walks every shape in every slide and concatenates the readable text
runs. Notes pages are included as well — Korean institutions often
embed reviewer comments there. Tables produce one cell per text-run.

Encrypted (MS-CFB) presentations surface as REQ-4051 to mirror the
docx / xlsx behaviour.
"""

from __future__ import annotations

import asyncio
import io
import logging
import zipfile

from app.extractors.fetcher import ExtractionError

logger = logging.getLogger(__name__)


def _is_encrypted_ooxml(data: bytes) -> bool:
    return data[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _shape_text(shape: object) -> list[str]:
    """Best-effort text harvest from a single python-pptx shape.

    Falls through quietly when the shape has no text frame (e.g. a
    picture or chart placeholder). Tables and grouped shapes recurse.
    """
    out: list[str] = []
    has_table = bool(getattr(shape, "has_table", False))
    if has_table:
        table = shape.table  # type: ignore[attr-defined]
        for row in table.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                if txt:
                    out.append(txt)
        return out

    # GroupShape — recurse into children. python-pptx ≥ 0.6 always
    # provides the MSO_SHAPE_TYPE enum; if the import unexpectedly fails
    # we just treat the shape as a leaf and continue without crashing.
    shape_type_attr = getattr(shape, "shape_type", None)
    if shape_type_attr is not None:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]

            is_group = shape_type_attr == MSO_SHAPE_TYPE.GROUP
        except Exception as e:
            logger.debug("pptx shape-type probe failed: %s", e)
            is_group = False
        if is_group:
            for child in getattr(shape, "shapes", []):
                out.extend(_shape_text(child))
            return out

    has_text = bool(getattr(shape, "has_text_frame", False))
    if not has_text:
        return out
    frame = shape.text_frame  # type: ignore[attr-defined]
    for para in frame.paragraphs:
        chunks = [run.text for run in para.runs if run.text]
        if chunks:
            out.append("".join(chunks))
    return out


def _extract_sync(data: bytes, filename: str) -> str:
    """Synchronous core run on a worker thread."""
    from pptx import Presentation  # type: ignore[import-not-found]

    if _is_encrypted_ooxml(data):
        raise ExtractionError("REQ-4051", filename=filename, detail="password-protected")

    try:
        prs = Presentation(io.BytesIO(data))
    except zipfile.BadZipFile as e:
        raise ExtractionError("REQ-4042", filename=filename, detail=f"bad zip: {e}") from e
    except Exception as e:
        raise ExtractionError("REQ-4042", filename=filename, detail=str(e)) from e

    parts: list[str] = []
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                parts.extend(_shape_text(shape))
            # Speaker notes — often hold reviewer comments with PII.
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                for para in notes_tf.paragraphs:
                    chunks = [run.text for run in para.runs if run.text]
                    if chunks:
                        parts.append("".join(chunks))
    except Exception as e:
        raise ExtractionError("REQ-4042", filename=filename, detail=f"walk failed: {e}") from e

    return "\n".join(parts)


async def extract_pptx(data: bytes, filename: str) -> str:
    """Extract concatenated text from a .pptx (OOXML PresentationML) file."""
    return await asyncio.to_thread(_extract_sync, data, filename)
